"""
PPT Generator Core
Creates professional PowerPoint presentations
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from typing import List, Dict, Any, Optional
import logging

logger = logging.getLogger(__name__)

class PPTGenerator:
    """PowerPoint presentation generator"""
    
    @staticmethod
    def _ensure_list(data):
        """Ensure data is a list, converting dict values or wrapping single items"""
        if isinstance(data, dict):
            return list(data.values()) if data else []
        if isinstance(data, list):
            return data
        if data is None:
            return []
        return [data]
    
    @staticmethod
    def _safe_str(data, max_len: int = 500) -> str:
        """Safely convert any data to string, handling dicts"""
        if isinstance(data, dict):
            # Try to find a text field, otherwise stringify
            for key in ['text', 'content', 'description', 'value', 'summary']:
                if key in data:
                    return str(data[key])[:max_len]
            return str(data)[:max_len]
        if data is None:
            return ""
        return str(data)[:max_len]
    
    COLOR_SCHEMES = {
        "professional": {
            "primary": RGBColor(30, 64, 175),
            "secondary": RGBColor(59, 130, 246),
            "accent": RGBColor(16, 185, 129),
            "text": RGBColor(31, 41, 55),
            "light": RGBColor(243, 244, 246),
            "white": RGBColor(255, 255, 255),
        },
        "minimal": {
            "primary": RGBColor(31, 41, 55),
            "secondary": RGBColor(107, 114, 128),
            "accent": RGBColor(99, 102, 241),
            "text": RGBColor(31, 41, 55),
            "light": RGBColor(249, 250, 251),
            "white": RGBColor(255, 255, 255),
        },
        "dark": {
            "primary": RGBColor(99, 102, 241),
            "secondary": RGBColor(139, 92, 246),
            "accent": RGBColor(236, 72, 153),
            "text": RGBColor(255, 255, 255),
            "light": RGBColor(55, 65, 81),
            "white": RGBColor(255, 255, 255),
        },
        "startup": {
            "primary": RGBColor(244, 63, 94),
            "secondary": RGBColor(251, 146, 60),
            "accent": RGBColor(234, 179, 8),
            "text": RGBColor(31, 41, 55),
            "light": RGBColor(255, 241, 242),
            "white": RGBColor(255, 255, 255),
        }
    }
    
    def create_presentation(
        self,
        project_name: str,
        project_description: Optional[str],
        steps: List[Dict],
        template: str,
        output_path: str
    ):
        """Create a complete presentation"""
        # Ensure steps is a list (handles dict input from API)
        steps = self._ensure_list(steps)
        
        colors = self.COLOR_SCHEMES.get(template, self.COLOR_SCHEMES["professional"])
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Title Slide
        self._add_title_slide(prs, project_name, project_description, colors)
        
        # Agenda
        self._add_agenda_slide(prs, steps, colors)
        
        # Content slides
        for step in steps:
            self._add_step_slide(prs, step, colors)
        
        # Closing
        self._add_closing_slide(prs, colors)
        
        prs.save(output_path)
        logger.info(f"Saved: {output_path}")
    
    def _add_title_slide(self, prs, title: str, subtitle: Optional[str], colors: Dict):
        """Create title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = colors["light"]
        bg.line.fill.background()
        
        # Accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            prs.slide_width, Inches(0.3)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors["primary"]
        bar.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
        frame = title_box.text_frame
        frame.text = title
        frame.paragraphs[0].font.size = Pt(48)
        frame.paragraphs[0].font.bold = True
        frame.paragraphs[0].font.color.rgb = colors["primary"]
        
        # Subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(10), Inches(1))
            sub_frame = sub_box.text_frame
            sub_frame.text = subtitle
            sub_frame.paragraphs[0].font.size = Pt(24)
            sub_frame.paragraphs[0].font.color.rgb = colors["secondary"]
    
    def _add_agenda_slide(self, prs, steps: List[Dict], colors: Dict):
        """Create agenda slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Header
        self._add_header_bar(slide, "Presentation Overview", colors, prs.slide_width)
        
        # Agenda items - ensure steps is a list
        steps = self._ensure_list(steps)
        y_pos = 1.8
        for i, step in enumerate(steps[:6]):
            step_name = step.get('stepName', f"Step {i+1}")
            
            # Number circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(1), Inches(y_pos), Inches(0.5), Inches(0.5)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = colors["primary"]
            circle.line.fill.background()
            
            # Number
            num_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.1), Inches(0.5), Inches(0.3))
            num_frame = num_box.text_frame
            num_frame.text = str(i + 1)
            num_frame.paragraphs[0].font.size = Pt(14)
            num_frame.paragraphs[0].font.bold = True
            num_frame.paragraphs[0].font.color.rgb = colors["white"]
            num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Step name
            name_box = slide.shapes.add_textbox(Inches(1.8), Inches(y_pos + 0.05), Inches(8), Inches(0.4))
            name_frame = name_box.text_frame
            name_frame.text = step_name
            name_frame.paragraphs[0].font.size = Pt(18)
            name_frame.paragraphs[0].font.color.rgb = colors["text"]
            
            y_pos += 0.8
    
    def _add_step_slide(self, prs, step: Dict, colors: Dict):
        """Add content slide for a step"""
        step_id = step.get('stepId', 0)
        step_name = step.get('stepName', 'Step')
        data = step.get('data', {})
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Header
        self._add_header_bar(slide, f"{step_id}. {step_name}", colors, prs.slide_width)
        
        # Content based on step type
        if step_id == 1:
            self._add_problem_slide(slide, data, colors)
        elif step_id == 2:
            self._add_vision_slide(slide, data, colors)
        elif step_id == 3:
            self._add_personas_slide(slide, data, colors)
        elif step_id == 4:
            self._add_questions_slide(slide, data, colors)
        elif step_id == 5:
            self._add_market_slide(slide, data, colors)
        elif step_id == 6:
            self._add_features_slide(slide, data, colors)
        elif step_id == 7:
            self._add_stories_slide(slide, data, colors)
        elif step_id == 8:
            self._add_roadmap_slide(slide, data, colors)
        elif step_id == 9:
            self._add_okrs_slide(slide, data, colors)
        else:
            self._add_generic_slide(slide, data, colors)
    
    def _add_header_bar(self, slide, title: str, colors: Dict, width):
        """Add header bar to slide"""
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            width, Inches(1)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = colors["primary"]
        header.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = colors["white"]
    
    def _add_problem_slide(self, slide, data: Dict, colors: Dict):
        """Problem reframe slide"""
        y_pos = 1.3
        
        if data.get('problemTitle'):
            problem = data['problemTitle']
            if isinstance(problem, dict):
                problem = str(problem)
            self._add_text_box(slide, "Problem:", problem, 
                             Inches(0.5), Inches(y_pos), colors)
            y_pos += 1.0
        
        if data.get('reframedProblem'):
            reframed = data['reframedProblem']
            if isinstance(reframed, dict):
                reframed = str(reframed)
            self._add_text_box(slide, "Reframed:", reframed,
                             Inches(0.5), Inches(y_pos), colors)
            y_pos += 1.2
        
        if data.get('rootCauses'):
            items = self._ensure_list(data['rootCauses'])
            if items:
                self._add_bullet_list(slide, "Root Causes:", items[:5],
                                    Inches(0.5), Inches(y_pos), colors)
    
    def _add_vision_slide(self, slide, data: Dict, colors: Dict):
        """Product vision slide"""
        if data.get('visionStatement'):
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(12), Inches(1.2)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = colors["light"]
            box.line.color.rgb = colors["primary"]
            
            frame = box.text_frame
            vision_text = data['visionStatement']
            if isinstance(vision_text, dict):
                vision_text = str(vision_text)
            frame.text = vision_text
            frame.paragraphs[0].font.size = Pt(20)
            frame.paragraphs[0].font.italic = True
            frame.paragraphs[0].font.color.rgb = colors["primary"]
            frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            frame.word_wrap = True
        
        y_pos = 2.8
        if data.get('elevatorPitch'):
            pitch = data['elevatorPitch']
            if isinstance(pitch, dict):
                pitch = str(pitch)
            self._add_text_box(slide, "Elevator Pitch:", pitch,
                             Inches(0.5), Inches(y_pos), colors)
            y_pos += 1.5
        
        if data.get('targetAudience'):
            audience = data['targetAudience']
            if isinstance(audience, dict):
                audience = str(audience)
            self._add_text_box(slide, "Target Audience:", audience,
                             Inches(0.5), Inches(y_pos), colors)
    
    def _add_personas_slide(self, slide, data: Dict, colors: Dict):
        """User personas slide"""
        personas = data.get('personas', [])
        if not isinstance(personas, list):
            personas = []
        
        personas = self._ensure_list(personas)
        if not personas:
            self._add_text_box(slide, "", "No persona data available",
                             Inches(0.5), Inches(1.5), colors)
            return
        
        card_width = Inches(4)
        gap = Inches(0.3)
        start_x = Inches(0.5)
        y_pos = Inches(1.3)
        
        for i, persona in enumerate(personas[:3]):
            x_pos = start_x + (i * (card_width + gap))
            self._add_persona_card(slide, persona, x_pos, y_pos, card_width, colors)
    
    def _add_persona_card(self, slide, persona: Dict, x, y, width, colors):
        """Create persona card"""
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, Inches(5.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = colors["light"]
        card.line.color.rgb = colors["secondary"]
        
        # Name
        name = persona.get('name', 'Unknown')
        role = persona.get('role', 'User')
        name_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2),
                                           width - Inches(0.4), Inches(0.5))
        name_frame = name_box.text_frame
        name_frame.text = f"{name} - {role}"
        name_frame.paragraphs[0].font.size = Pt(16)
        name_frame.paragraphs[0].font.bold = True
        name_frame.paragraphs[0].font.color.rgb = colors["primary"]
        
        # Bio
        bio = self._safe_str(persona.get('bio', ''), 150)
        if len(bio) > 150:
            bio = bio[:150] + "..."
        self._add_wrapped_text(slide, bio, x + Inches(0.2), y + Inches(0.8),
                              width - Inches(0.4), Pt(10), colors["text"])
    
    def _add_questions_slide(self, slide, data: Dict, colors: Dict):
        """Q&A slide"""
        questions = data.get('questions', [])
        if not isinstance(questions, list):
            questions = []
        
        questions = self._ensure_list(questions)
        if not questions:
            self._add_text_box(slide, "", "No questions available",
                             Inches(0.5), Inches(1.5), colors)
            return
        
        y_pos = 1.3
        for i, q in enumerate(questions[:6]):
            question = self._safe_str(q.get('question', ''), 80)
            answer = self._safe_str(q.get('aiAnswer', q.get('userAnswer', '')), 100)
            
            # Question
            q_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12), Inches(0.4))
            q_frame = q_box.text_frame
            q_frame.text = f"Q{i+1}: {question}"
            q_frame.paragraphs[0].font.size = Pt(12)
            q_frame.paragraphs[0].font.bold = True
            q_frame.paragraphs[0].font.color.rgb = colors["primary"]
            
            # Answer
            if answer:
                a_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.4), Inches(11.5), Inches(0.6))
                a_frame = a_box.text_frame
                a_frame.text = answer
                a_frame.paragraphs[0].font.size = Pt(10)
                a_frame.paragraphs[0].font.color.rgb = colors["text"]
                a_frame.word_wrap = True
            
            y_pos += 1.2
    
    def _add_market_slide(self, slide, data: Dict, colors: Dict):
        """Market analysis slide"""
        market_overview = data.get('marketOverview')
        if market_overview:
            # Handle both string and dict input
            if isinstance(market_overview, dict):
                overview = str(market_overview)[:200]
            else:
                overview = str(market_overview)[:200]
            self._add_text_box(slide, "Market Overview:", overview,
                             Inches(0.5), Inches(1.3), colors)
        
        competitors = self._ensure_list(data.get('competitors', []))
        if competitors:
            self._add_competitor_table(slide, competitors[:4], colors)
    
    def _add_competitor_table(self, slide, competitors: List[Dict], colors: Dict):
        """Add competitor comparison table"""
        rows = len(competitors) + 1
        table = slide.shapes.add_table(rows, 3, Inches(0.5), Inches(3),
                                      Inches(12), Inches(0.5 * rows)).table
        
        headers = ["Competitor", "Strengths", "Weaknesses"]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors["primary"]
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = colors["white"]
        
        for i, comp in enumerate(competitors):
            row = i + 1
            table.cell(row, 0).text = comp.get('name', 'Unknown')
            
            strengths = self._ensure_list(comp.get('strengths', []))
            if strengths:
                table.cell(row, 1).text = "\n".join([f"• {self._safe_str(s, 40)}" for s in strengths[:2]])
            
            weaknesses = self._ensure_list(comp.get('weaknesses', []))
            if weaknesses:
                table.cell(row, 2).text = "\n".join([f"• {self._safe_str(w, 40)}" for w in weaknesses[:2]])
    
    def _add_features_slide(self, slide, data: Dict, colors: Dict):
        """PRD/Features slide"""
        features = data.get('features', [])
        if not isinstance(features, list):
            features = []
        
        features = self._ensure_list(features)
        if not features:
            self._add_text_box(slide, "PRD Content", str(data)[:300],
                             Inches(0.5), Inches(1.5), colors)
            return
        
        y_pos = 1.5
        for feature in features[:8]:
            name = self._safe_str(feature.get('name', feature.get('title', str(feature))), 60)
            priority = self._safe_str(feature.get('priority', 'Medium'), 20)
            
            # Priority badge
            badge = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(0.8), Inches(0.3)
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = self._get_priority_color(priority, colors)
            badge.line.fill.background()
            
            # Feature name
            feature_box = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos), Inches(11), Inches(0.3))
            feature_frame = feature_box.text_frame
            feature_frame.text = name
            feature_frame.paragraphs[0].font.size = Pt(12)
            feature_frame.paragraphs[0].font.color.rgb = colors["text"]
            
            y_pos += 0.45
    
    def _get_priority_color(self, priority: str, colors: Dict):
        """Get color for priority level"""
        p = priority.lower()
        if p == 'high':
            return RGBColor(239, 68, 68)
        elif p == 'medium':
            return RGBColor(245, 158, 11)
        return RGBColor(16, 185, 129)
    
    def _add_stories_slide(self, slide, data: Dict, colors: Dict):
        """User stories slide"""
        stories = data.get('stories', [])
        if not isinstance(stories, list):
            stories = []
        
        stories = self._ensure_list(stories)
        if not stories:
            self._add_text_box(slide, "", "No user stories available",
                             Inches(0.5), Inches(1.5), colors)
            return
        
        rows = min(len(stories) + 1, 7)
        table = slide.shapes.add_table(rows, 4, Inches(0.5), Inches(1.5),
                                      Inches(12), Inches(0.5 * rows)).table
        
        headers = ["ID", "Story", "Priority", "RICE"]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors["primary"]
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = colors["white"]
        
        table.columns[0].width = Inches(0.8)
        table.columns[1].width = Inches(8)
        table.columns[2].width = Inches(1.2)
        table.columns[3].width = Inches(1)
        
        for i, story in enumerate(stories[:6]):
            row = i + 1
            table.cell(row, 0).text = str(story.get('id', i + 1))
            
            desc = story.get('description', '')
            if not desc and story.get('asA'):
                desc = f"As a {story.get('asA')}, I want {story.get('iWant', '...')}"
            table.cell(row, 1).text = self._safe_str(desc, 80)
            
            table.cell(row, 2).text = str(story.get('priority', '-'))
            
            rice = story.get('riceScore', '-')
            table.cell(row, 3).text = str(rice)[:4]
    
    def _add_roadmap_slide(self, slide, data: Dict, colors: Dict):
        """Roadmap slide"""
        phases = data.get('phases', [])
        if not isinstance(phases, list):
            phases = []
        
        phases = self._ensure_list(phases)
        if not phases:
            self._add_text_box(slide, "", "No roadmap data available",
                             Inches(0.5), Inches(1.5), colors)
            return
        
        y_pos = 1.5
        for i, phase in enumerate(phases[:4]):
            # Card
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(12), Inches(1.2)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = colors["light"]
            card.line.color.rgb = colors["secondary"]
            
            # Phase number
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(0.8), Inches(y_pos + 0.3), Inches(0.6), Inches(0.6)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = colors["primary"]
            circle.line.fill.background()
            
            num_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.4), Inches(0.6), Inches(0.4))
            num_frame = num_box.text_frame
            num_frame.text = str(i + 1)
            num_frame.paragraphs[0].font.size = Pt(18)
            num_frame.paragraphs[0].font.bold = True
            num_frame.paragraphs[0].font.color.rgb = colors["white"]
            num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Phase name
            name_box = slide.shapes.add_textbox(Inches(1.7), Inches(y_pos + 0.2), Inches(4), Inches(0.4))
            name_frame = name_box.text_frame
            name_frame.text = phase.get('name', f'Phase {i + 1}')
            name_frame.paragraphs[0].font.size = Pt(16)
            name_frame.paragraphs[0].font.bold = True
            name_frame.paragraphs[0].font.color.rgb = colors["primary"]
            
            y_pos += 1.4
    
    def _add_okrs_slide(self, slide, data: Dict, colors: Dict):
        """OKRs slide"""
        okrs = [data.get('okr1'), data.get('okr2'), data.get('okr3')]
        okrs = [okr for okr in okrs if okr]
        
        if not okrs:
            self._add_text_box(slide, "", "No OKR data available",
                             Inches(0.5), Inches(1.5), colors)
            return
        
        if data.get('northStarDefinition'):
            self._add_info_card(slide, "North Star", data['northStarDefinition'],
                              Inches(0.5), Inches(1.4), Inches(12), Inches(0.8), colors)
        
        y_pos = 2.4
        for okr in okrs[:3]:
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(12), Inches(1.3)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = colors["light"]
            card.line.color.rgb = colors["primary"]
            
            obj = okr.get('objective', '') if isinstance(okr, dict) else str(okr)[:80]
            obj_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.1), Inches(11.5), Inches(0.4))
            obj_frame = obj_box.text_frame
            obj_frame.text = obj
            obj_frame.paragraphs[0].font.size = Pt(14)
            obj_frame.paragraphs[0].font.bold = True
            obj_frame.paragraphs[0].font.color.rgb = colors["primary"]
            
            y_pos += 1.5
    
    def _add_generic_slide(self, slide, data: Dict, colors: Dict):
        """Generic content slide"""
        content = str(data)[:500]
        self._add_text_box(slide, "Content", content,
                          Inches(0.5), Inches(1.5), colors)
    
    def _add_closing_slide(self, prs, colors: Dict):
        """Thank you slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = colors["primary"]
        bg.line.fill.background()
        
        text_box = slide.shapes.add_textbox(Inches(0), Inches(3), prs.slide_width, Inches(1.5))
        frame = text_box.text_frame
        frame.text = "Thank You"
        frame.paragraphs[0].font.size = Pt(60)
        frame.paragraphs[0].font.bold = True
        frame.paragraphs[0].font.color.rgb = colors["white"]
        frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _add_text_box(self, slide, label: str, content, x, y, colors: Dict):
        """Add labeled text box"""
        if label:
            label_box = slide.shapes.add_textbox(x, y, Inches(12), Inches(0.3))
            label_frame = label_box.text_frame
            label_frame.text = label
            label_frame.paragraphs[0].font.size = Pt(11)
            label_frame.paragraphs[0].font.bold = True
            label_frame.paragraphs[0].font.color.rgb = colors["primary"]
            y += Inches(0.35)
        
        # Convert content to string safely
        content = self._safe_str(content, 300)
        
        content_box = slide.shapes.add_textbox(x, y, Inches(12), Inches(2))
        content_frame = content_box.text_frame
        content_frame.text = content
        content_frame.paragraphs[0].font.size = Pt(12)
        content_frame.paragraphs[0].font.color.rgb = colors["text"]
        content_frame.word_wrap = True
    
    def _add_bullet_list(self, slide, label: str, items: List, x, y, colors: Dict):
        """Add bullet list"""
        items = self._ensure_list(items)
        if not items:
            return
        
        if label:
            label_box = slide.shapes.add_textbox(x, y, Inches(12), Inches(0.3))
            label_frame = label_box.text_frame
            label_frame.text = label
            label_frame.paragraphs[0].font.size = Pt(11)
            label_frame.paragraphs[0].font.bold = True
            label_frame.paragraphs[0].font.color.rgb = colors["primary"]
            y += Inches(0.35)
        
        bullet_text = "\n".join([f"• {self._safe_str(item, 60)}" for item in items[:5]])
        bullet_box = slide.shapes.add_textbox(x, y, Inches(12), Inches(2))
        bullet_frame = bullet_box.text_frame
        bullet_frame.text = bullet_text
        bullet_frame.paragraphs[0].font.size = Pt(11)
        bullet_frame.paragraphs[0].font.color.rgb = colors["text"]
    
    def _add_wrapped_text(self, slide, text: str, x, y, width, font_size, color):
        """Add wrapped text"""
        box = slide.shapes.add_textbox(x, y, width, Inches(1))
        frame = box.text_frame
        frame.text = text
        frame.paragraphs[0].font.size = font_size
        frame.paragraphs[0].font.color.rgb = color
        frame.word_wrap = True
    
    def _add_info_card(self, slide, title: str, content, x, y, width, height, colors: Dict):
        """Add info card"""
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = colors["light"]
        card.line.color.rgb = colors["primary"]
        
        title_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.05),
                                            width - Inches(0.2), Inches(0.3))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(11)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = colors["primary"]
        
        # Convert content to string safely
        content = self._safe_str(content, 200)
        
        content_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.35),
                                              width - Inches(0.2), height - Inches(0.4))
        content_frame = content_box.text_frame
        content_frame.text = content
        content_frame.paragraphs[0].font.size = Pt(10)
        content_frame.word_wrap = True
